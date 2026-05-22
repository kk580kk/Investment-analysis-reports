from __future__ import annotations

import hashlib
import json
import sys
from typing import cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import XmlPart, PartFactory
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, _nsmap
from pptx.oxml.xmlchemy import BaseOxmlElement
from pptx.package import Package
from pptx.util import Inches, Pt


# ═══════════════════════════════════════════════════════════════
#  1. 注册 vt 命名空间
# ═══════════════════════════════════════════════════════════════

_nsmap.update({
    "vt": "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes",
})


# ═══════════════════════════════════════════════════════════════
#  2. OxmlElement 类（原 ppt_extend/oxml/customprops.py）
# ═══════════════════════════════════════════════════════════════

class CT_Lpwstr(BaseOxmlElement):
    """<vt:lpwstr> 元素。"""

    @classmethod
    def new(cls, value: str) -> CT_Lpwstr:
        xml = "<vt:lpwstr %s>%s</vt:lpwstr>" % (nsdecls("vt"), value)
        return cast(CT_Lpwstr, parse_xml(xml))


class CT_Property(BaseOxmlElement):
    """单条 custom property 元素。"""

    @classmethod
    def new(cls, key: str, value: str, fmtid: str, pid: int) -> CT_Property:
        xml = '<property fmtid="{%s}" pid="%d" name="%s"/>' % (fmtid, pid, key)
        prop = cast(CT_Property, parse_xml(xml))
        prop.append(CT_Lpwstr.new(value))
        return prop


class CT_CustomProperties(BaseOxmlElement):
    """<Properties> 根元素。"""

    @classmethod
    def new(cls) -> CT_CustomProperties:
        xml = (
            '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/'
            '2006/custom-properties" %s/>\n' % nsdecls("vt")
        )
        return cast(CT_CustomProperties, parse_xml(xml))

    def add_property(self, prop: CT_Property) -> CT_Property:
        self.append(prop)
        return prop


# ═══════════════════════════════════════════════════════════════
#  3. CustomProperties 包装
# ═══════════════════════════════════════════════════════════════

class CustomProperties:
    """对 CT_CustomProperties 的高层封装。"""

    def __init__(self, element: CT_CustomProperties):
        self._element = element

    def add_property(self, key: str, value: str, fmtid: str, pid: int):
        prop = CT_Property.new(key, value, fmtid, pid)
        self._element.append(prop)


# ═══════════════════════════════════════════════════════════════
#  4. CustomPropertiesPart
# ═══════════════════════════════════════════════════════════════

class CustomPropertiesPart(XmlPart):
    """对应 /docProps/custom.xml 部件。"""

    _element: CT_CustomProperties

    @classmethod
    def default(cls, package: Package) -> CustomPropertiesPart:
        partname = PackURI("/docProps/custom.xml")
        content_type = CT.OFC_CUSTOM_PROPERTIES
        element = CT_CustomProperties.new()
        return cls(partname, content_type, package, element)

    @property
    def custom_properties(self) -> CustomProperties:
        return CustomProperties(self._element)

    @property
    def next_id(self) -> int:
        id_str_lst = self._element.xpath("//@pid")
        used_ids = [int(s) for s in id_str_lst if s.isdigit()]
        return max(used_ids) + 1 if used_ids else 2


# 注册 PartFactory（原 ppt_extend/__init__.py）
PartFactory.part_type_for[CT.OFC_CUSTOM_PROPERTIES] = CustomPropertiesPart


# ═══════════════════════════════════════════════════════════════
#  5. PresentationExtend
# ═══════════════════════════════════════════════════════════════

class PresentationExtend:
    """为 python-pptx Presentation 扩展 custom properties 访问能力。"""

    def __init__(self, prs: Presentation):
        self._prs = prs

    @property
    def custom_properties(self) -> CustomProperties:
        return self.custom_properties_part.custom_properties

    @property
    def custom_properties_part(self) -> CustomPropertiesPart:
        package = self._prs._part.package
        try:
            return package.part_related_by(RT.CUSTOM_PROPERTIES)
        except KeyError:
            part = CustomPropertiesPart.default(package)
            package.relate_to(part, RT.CUSTOM_PROPERTIES)
            return part


# ═══════════════════════════════════════════════════════════════
#  6. AIGC 签名生成 + 文本提取 + 水印
# ═══════════════════════════════════════════════════════════════

_CUSTOM_PROPERTY_FMTID = "D5CDD505-2E9C-101B-9397-08002B2CF9AE"


def generate_sha256(text: str) -> str:
    """计算字符串的 SHA256 哈希值。"""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def get_aigc_signature(content: str, request_id: str = "") -> str:
    """生成 AIGC 隐式标识 JSON 字符串。"""
    content_id = generate_sha256(content)[:16]
    produce_id = f"voiceassistant-{content_id}"
    return json.dumps(
        {
            "Label": "1",
            "ContentProducer": "001191320114777023172010000",
            "ProduceID": produce_id,
            "ReservedCode1": "",
            "ContentPropagator": "001191320114777023172010000",
            "PropagateID": produce_id,
        },
        ensure_ascii=False,
    )


def extract_all_text(prs: Presentation, request_id: str = "") -> str:
    """从 PPT 中提取所有文字内容。"""
    try:
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            t = run.text.strip()
                            if t:
                                texts.append(t)
        return "\n".join(texts)
    except Exception as e:
        print(f"{request_id}, 读取 PPT 文件时出错: {e}")
        return ""


def _add_watermark_first_slide(prs: Presentation, request_id: str = ""):
    """在第一页底部右侧添加「内容由AI生成」文本标识。"""
    if not prs.slides:
        print(f"request_id={request_id}, no slides found")
        return

    slide = prs.slides[0]
    left = prs.slide_width - Inches(3.2)
    top = prs.slide_height - Inches(0.5)

    txbox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.3))
    tf = txbox.text_frame
    tf.text = "内容由AI生成"

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(128, 128, 128)
    run.font.name = "微软雅黑"

    print(f"request_id={request_id}, added AI watermark to first slide")


def add_aigc_mark(prs: Presentation, aigc_signature: str, request_id: str = "") -> None:
    """向已打开的 Presentation 添加可见水印 + 隐式 custom property。"""
    try:
        _add_watermark_first_slide(prs, request_id)
        prs_ex = PresentationExtend(prs)
        part = prs_ex.custom_properties_part
        part.custom_properties.add_property(
            "AIGC", aigc_signature, _CUSTOM_PROPERTY_FMTID, part.next_id
        )
    except Exception as e:
        print(f"request_id={request_id}, add AIGC mark error: {e}")


def add_aigc_mark_to_pptx(
    input_path: str,
    output_path: str,
    aigc_signature: str,
    request_id: str = "",
) -> None:
    """一步到位：读取 → 加标记 → 保存。"""
    prs = Presentation(input_path)
    add_aigc_mark(prs, aigc_signature, request_id)
    prs.save(output_path)


# ═══════════════════════════════════════════════════════════════
#  7. 便捷入口：自动生成签名并打标
# ═══════════════════════════════════════════════════════════════

def mark_pptx(input_path: str, output_path: str | None = None, request_id: str = ""):
    """最简调用：只需传入文件路径即可完成全部标记。

    Args:
        input_path:  源 .pptx 文件
        output_path: 输出路径，默认覆盖源文件
        request_id:  可选追踪 ID
    """
    if output_path is None:
        output_path = input_path
    prs = Presentation(input_path)
    all_text = extract_all_text(prs, request_id)
    signature = get_aigc_signature(all_text, request_id)
    add_aigc_mark(prs, signature, request_id)
    prs.save(output_path)
    print(f"[PPTX] AIGC mark added → {output_path}")


# ═══════════════════════════════════════════════════════════════
#  CLI
# ═══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python pptx_aigc_mark.py <input.pptx> [output.pptx]")
        sys.exit(1)
    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else src
    mark_pptx(src, dst)