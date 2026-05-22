from pptx import Presentation


def preprocess_pptx_for_pdf(input_path, output_path):
    """预处理PPTX用于PDF转换：当单页无文字shape数量超过20个时，删除该页所有无文字shape

    某些PPTX文件中的矢量图会被拆分成大量丝状小shape，导致soffice转换时超时。
    此预处理通过删除高密度空shape页中的无文字shape来避免超时。

    Args:
        input_path: 输入pptx文件路径
        output_path: 输出pptx文件路径

    Returns:
        bool: 是否成功处理
    """
    try:
        prs = Presentation(input_path)
        total_removed = 0
        modified_slides = 0

        for slide in prs.slides:
            empty_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    if not shape.text_frame.text.strip():
                        empty_shapes.append(shape)

            if len(empty_shapes) > 20:
                for shape in reversed(empty_shapes):
                    sp = shape._element
                    sp.getparent().remove(sp)
                    total_removed += 1
                modified_slides += 1

        prs.save(output_path)
        if total_removed > 0:
            print(
                f"Preprocessed PPTX for PDF: removed {total_removed} empty shapes "
                f"from {modified_slides} slide(s)"
            )
        return True
    except Exception as e:
        print(f"Failed to preprocess PPTX {input_path}: {e}")
        return False
